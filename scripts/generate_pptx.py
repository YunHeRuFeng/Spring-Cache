import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# =========================
# Global theme (premium-minimal)
# =========================
BG = RGBColor(0x0A, 0x0E, 0x1A)           # #0A0E1A
SURFACE = RGBColor(0x12, 0x18, 0x28)      # dark card
SURFACE_2 = RGBColor(0x16, 0x1F, 0x33)    # alt dark
BORDER = RGBColor(0x2A, 0x36, 0x4F)
TEXT_MAIN = RGBColor(0xE0, 0xE6, 0xF0)    # #E0E6F0
TEXT_SUB = RGBColor(0x94, 0xA3, 0xB8)     # #94A3B8
TEXT_DIM = RGBColor(0x6B, 0x7E, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

ACCENT_CYAN = RGBColor(0x38, 0xBD, 0xF8)
ACCENT_INDIGO = RGBColor(0x81, 0x8C, 0xF8)
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
ACCENT_PINK = RGBColor(0xF4, 0x72, 0xB6)
ACCENT_GREEN = RGBColor(0x34, 0xD3, 0x99)
ACCENT_RED = RGBColor(0xF8, 0x71, 0x71)
ACCENT_ORANGE = RGBColor(0xFB, 0x92, 0x3C)
ACCENT_YELLOW = RGBColor(0xFB, 0xBF, 0x24)

SLIDE_W = 13.33
SLIDE_H = 7.5
TOTAL = 7


# =========================
# Core utils
# =========================
def set_run(run, size=14, bold=False, color=TEXT_MAIN, font="Microsoft YaHei"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=TEXT_MAIN, align=PP_ALIGN.LEFT, font="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_run(r, size=size, bold=bold, color=color, font=font)
    return tb


def add_multiline(slide, x, y, w, h, lines, size=12, color=TEXT_SUB):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"• {line}"
        set_run(r, size=size, color=color)
    return tb


def bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    # subtle corner glow (very light, avoid noisy look)
    g1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.6), Inches(-1.2), Inches(4.8), Inches(3.6))
    g1.fill.solid()
    g1.fill.fore_color.rgb = ACCENT_INDIGO
    g1.fill.transparency = 0.92
    g1.line.fill.background()

    g2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(5.0), Inches(4.2), Inches(3.1))
    g2.fill.solid()
    g2.fill.fore_color.rgb = ACCENT_CYAN
    g2.fill.transparency = 0.93
    g2.line.fill.background()


def top_bar(slide, title, source=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.28), Inches(12.25), Inches(0.62))
    bar.fill.solid()
    bar.fill.fore_color.rgb = SURFACE
    bar.line.color.rgb = BORDER
    bar.line.width = Pt(1)

    # left accent
    ac = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.28), Inches(0.08), Inches(0.62))
    ac.fill.solid()
    ac.fill.fore_color.rgb = ACCENT_CYAN
    ac.line.fill.background()

    add_text(slide, 0.78, 0.34, 8.3, 0.46, title, size=21, bold=True, color=WHITE)
    if source:
        add_text(slide, 8.8, 0.38, 3.9, 0.34, source, size=10, color=TEXT_DIM, align=PP_ALIGN.RIGHT, font="Arial")


def footer(slide, idx):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(7.14), Inches(12.25), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x1D, 0x27, 0x3B)
    line.line.fill.background()

    progress_w = 12.25 * idx / TOTAL
    p = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(7.14), Inches(progress_w), Inches(0.06))
    p.fill.solid()
    p.fill.fore_color.rgb = ACCENT_CYAN
    p.line.fill.background()

    add_text(slide, 11.25, 6.94, 1.55, 0.2, f"{idx:02d}/{TOTAL:02d}", size=9.5, color=TEXT_DIM, align=PP_ALIGN.RIGHT, font="Arial")


def card(slide, x, y, w, h, title, body=None, accent=ACCENT_CYAN, surface=SURFACE):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid()
    c.fill.fore_color.rgb = surface
    c.line.color.rgb = BORDER
    c.line.width = Pt(1)

    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()

    add_text(slide, x + 0.16, y + 0.08, w - 0.24, 0.26, title, size=13.5, bold=True, color=WHITE)
    if body:
        add_text(slide, x + 0.16, y + 0.36, w - 0.24, h - 0.42, body, size=11.2, color=TEXT_SUB)
    return c


def chip(slide, x, y, w, h, text, color):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    add_text(slide, x, y, w, h, text, size=11.2, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Arial")


def arrow_right(slide, x, y, w=0.3, h=0.18, color=TEXT_SUB):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()


def arrow_down(slide, x, y, w=0.28, h=0.26, color=TEXT_SUB):
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()


def connect(slide, x1, y1, x2, y2, color=TEXT_SUB):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(1)
    return ln


# =========================
# Slide builders
# =========================
def slide1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)

    # refined hero block
    hero = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.06), Inches(11.4), Inches(4.95))
    hero.fill.solid()
    hero.fill.fore_color.rgb = SURFACE
    hero.fill.transparency = 0.06
    hero.line.color.rgb = BORDER
    hero.line.width = Pt(1.1)

    add_text(s, 1.25, 1.58, 10.8, 0.95, "Local Cache Framework", size=50, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Arial")
    add_text(s, 1.25, 2.52, 10.8, 0.5, "基于 Spring Boot 的本地缓存框架", size=20, color=TEXT_MAIN, align=PP_ALIGN.CENTER)

    tags = ["Java 17", "Spring Boot 3.3", "Thymeleaf", "ConcurrentHashMap", "ReentrantLock"]
    x = 1.25
    for t in tags:
        card(s, x, 3.36, 2.0, 0.54, t, accent=ACCENT_CYAN, surface=SURFACE_2)
        x += 2.15

    add_text(s, 1.25, 4.34, 10.8, 0.45, "Strategy / Template Method / Facade", size=18, bold=True, color=ACCENT_PURPLE, align=PP_ALIGN.CENTER, font="Arial")
    add_text(s, 1.25, 5.01, 10.8, 0.32, "Designing pluggable eviction policies with clear service boundaries", size=11, color=TEXT_DIM, align=PP_ALIGN.CENTER, font="Arial")

    footer(s, 1)


def slide2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    top_bar(s, "系统架构全景图（精华版）", "extracted from 01-architecture.html")

    layers = [
        ("CLIENT", "Dashboard 页面 + REST Client", ACCENT_CYAN),
        ("WEB", "DashboardController / CacheRestController", ACCENT_INDIGO),
        ("SERVICE", "CacheManagerService(门面) + ProductRepository", ACCENT_PURPLE),
        ("CACHE CORE", "Cache / LocalCache / EvictionPolicy / LRU-LFU-FIFO", ACCENT_PINK),
        ("DATA", "ConcurrentHashMap + 双向链表 + 频率桶", ACCENT_GREEN),
    ]

    y = 1.1
    for i, (title, body, ac) in enumerate(layers):
        card(s, 0.92, y, 11.96, 0.9, title, body, accent=ac, surface=SURFACE)
        if i < len(layers) - 1:
            arrow_down(s, 6.53, y + 0.92, 0.24, 0.18, TEXT_SUB)
        y += 1.03

    card(
        s,
        0.92,
        6.35,
        11.96,
        0.62,
        "演讲重点",
        "一句话：控制层接请求，服务层做编排，核心层可插拔策略，底层结构保证 O(1) 级操作",
        accent=ACCENT_CYAN,
        surface=SURFACE_2,
    )
    footer(s, 2)


def slide3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    top_bar(s, "请求生命周期（精华版）", "extracted from 02-request-lifecycle.html")

    chip(s, 0.9, 1.25, 2.2, 0.52, "1. HTTP 请求", ACCENT_CYAN)
    arrow_right(s, 3.18, 1.43, 0.26, 0.14, TEXT_SUB)
    chip(s, 3.52, 1.25, 2.8, 0.52, "2. cache.lookup", ACCENT_INDIGO)
    arrow_right(s, 6.40, 1.43, 0.26, 0.14, TEXT_SUB)
    chip(s, 6.74, 1.25, 2.8, 0.52, "3. 判断TTL/存在性", ACCENT_PURPLE)

    card(s, 0.88, 2.15, 3.95, 2.5, "HIT 分支", "直接返回缓存\nrecordAccess\nhitCount++", accent=ACCENT_GREEN, surface=SURFACE)
    card(s, 4.69, 2.15, 3.95, 2.5, "MISS 分支", "回源 ProductRepository\n写入 cache.put\nmissCount++", accent=ACCENT_ORANGE, surface=SURFACE)
    card(s, 8.50, 2.15, 3.38, 2.5, "EXPIRED 分支", "移除旧条目\n回源并重建缓存\nexpiredCount++", accent=ACCENT_YELLOW, surface=SURFACE)

    chip(s, 1.18, 4.95, 2.05, 0.46, "source = HIT", ACCENT_GREEN)
    chip(s, 5.02, 4.95, 2.05, 0.46, "source = MISS", ACCENT_ORANGE)
    chip(s, 8.84, 4.95, 2.45, 0.46, "source = EXPIRED", ACCENT_YELLOW)

    card(
        s,
        0.88,
        5.65,
        11.98,
        1.22,
        "容量满时的公共动作",
        "size ≥ capacity 时触发 evictKey()；不同策略仅改变 victim 的选择规则（LRU / LFU / FIFO）",
        accent=ACCENT_RED,
        surface=SURFACE_2,
    )
    footer(s, 3)


def slide4(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    top_bar(s, "LRU 策略原理（精华版）", "extracted from 03-lru-policy.html")

    card(
        s,
        0.88,
        1.05,
        11.98,
        0.86,
        "核心定义",
        "最近访问的元素最不容易被淘汰；访问命中会改变顺序",
        accent=ACCENT_CYAN,
    )

    add_text(s, 1.0, 2.28, 11.6, 0.5, "访问前：head → K1 → K2 → K3 → K4 → tail", size=20, color=WHITE, align=PP_ALIGN.CENTER, font="Arial")
    add_text(s, 1.0, 2.92, 11.6, 0.5, "访问 K2 后：head → K1 → K3 → K4 → K2 → tail", size=20, color=ACCENT_GREEN, align=PP_ALIGN.CENTER, font="Arial")

    card(s, 0.88, 3.72, 3.83, 1.72, "数据结构", "HashMap + 双向链表\nnodeIndex O(1) 定位", accent=ACCENT_PURPLE)
    card(s, 4.94, 3.72, 3.83, 1.72, "关键动作", "onAccess -> moveToTail\nevictKey -> head.next", accent=ACCENT_INDIGO)
    card(s, 8.99, 3.72, 2.87, 1.72, "复杂度", "查询/更新/淘汰\n均可到 O(1)", accent=ACCENT_GREEN)

    card(
        s,
        0.88,
        5.70,
        11.98,
        1.10,
        "演讲口诀",
        "写入尾插、访问移尾、淘汰弹头",
        accent=ACCENT_CYAN,
        surface=SURFACE_2,
    )
    footer(s, 4)


def slide5(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    top_bar(s, "LFU 策略原理（精华版）", "extracted from 04-lfu-policy.html")

    card(s, 0.88, 1.05, 5.75, 1.72, "结构一：frequencies", "记录 key -> 当前访问频次\n如 K1->1, K2->2, K3->3", accent=ACCENT_PURPLE)
    card(s, 7.11, 1.05, 4.75, 1.72, "结构二：keysByFrequency", "记录 freq -> keys 集合\n同频次用 LinkedHashSet 保序", accent=ACCENT_GREEN)

    card(s, 0.88, 3.02, 11.98, 0.98, "核心指针", "minFrequency 始终指向当前最小频次桶；淘汰时先看它", accent=ACCENT_ORANGE)

    chip(s, 0.92, 4.35, 2.75, 0.75, "1) 新键入 freq=1", ACCENT_INDIGO)
    chip(s, 3.92, 4.35, 2.75, 0.75, "2) 命中后频次+1", ACCENT_INDIGO)
    chip(s, 6.92, 4.35, 2.75, 0.75, "3) 空桶则更新min", ACCENT_INDIGO)
    chip(s, 9.92, 4.35, 2.00, 0.75, "4) 按min淘汰", ACCENT_INDIGO)

    connect(s, 3.72, 4.72, 3.92, 4.72, TEXT_DIM)
    connect(s, 6.72, 4.72, 6.92, 4.72, TEXT_DIM)
    connect(s, 9.72, 4.72, 9.92, 4.72, TEXT_DIM)

    card(
        s,
        0.88,
        5.45,
        11.98,
        1.35,
        "演讲结论",
        "LFU 先比较访问频次，再比较同频次先后顺序；适合高频热点稳定的场景",
        accent=ACCENT_PINK,
        surface=SURFACE_2,
    )
    footer(s, 5)


def slide6(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    top_bar(s, "FIFO 策略原理（精华版）", "extracted from 05-fifo-policy.html")

    card(s, 0.88, 1.05, 11.98, 0.9, "核心定义", "只按进入时间淘汰；访问不会改变顺序", accent=ACCENT_CYAN)

    add_text(s, 1.0, 2.28, 11.6, 0.5, "当前顺序：head → K1 → K2 → K3 → K4 → tail", size=19, color=WHITE, align=PP_ALIGN.CENTER, font="Arial")
    add_text(s, 1.0, 2.9, 11.6, 0.5, "访问 K2/K4 后：顺序不变（onAccess = no-op）", size=18, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER, font="Arial")
    add_text(s, 1.0, 3.5, 11.6, 0.5, "容量满后：淘汰 K1，插入 K5 → head → K2 → K3 → K4 → K5 → tail", size=17, color=ACCENT_GREEN, align=PP_ALIGN.CENTER, font="Arial")

    card(s, 0.88, 4.35, 3.83, 1.72, "写入规则", "append 到尾部\n保持到达顺序", accent=ACCENT_CYAN)
    card(s, 4.94, 4.35, 3.83, 1.72, "访问规则", "只读不改序\n实现最稳定", accent=ACCENT_INDIGO)
    card(s, 8.99, 4.35, 2.87, 1.72, "淘汰规则", "evictKey=head.next\n最早进入先出", accent=ACCENT_RED)

    footer(s, 6)


def table_cell(slide, x, y, w, h, text, fill, bold=False):
    cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.line.color.rgb = BORDER
    cell.line.width = Pt(0.9)

    add_text(
        slide,
        x + 0.04,
        y + 0.02,
        w - 0.08,
        h - 0.04,
        text,
        size=10.6,
        bold=bold,
        color=WHITE if bold else TEXT_MAIN,
        align=PP_ALIGN.CENTER,
    )


def slide7(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    top_bar(s, "三种策略对比总结", "LRU / LFU / FIFO")

    cols = [2.2, 3.35, 3.35, 3.35]
    x0, y0, rh = 0.62, 1.0, 0.67

    headers = ["维度", "LRU", "LFU", "FIFO"]
    x = x0
    for i, h in enumerate(headers):
        table_cell(s, x, y0, cols[i], rh, h, SURFACE_2, bold=True)
        x += cols[i]

    rows = [
        ("数据结构", "HashMap + 双向链表", "HashMap + 频次桶", "HashMap + 双向链表"),
        ("写入行为", "尾插", "freq=1 入桶", "尾插"),
        ("访问行为", "移尾", "频次 +1", "不变"),
        ("淘汰对象", "最久未访问", "最少访问（同频按序）", "最早写入"),
        ("时间复杂度", "O(1)", "均摊 O(1)", "O(1)"),
    ]

    y = y0 + rh
    for row in rows:
        x = x0
        table_cell(s, x, y, cols[0], rh, row[0], SURFACE, bold=True)
        x += cols[0]
        table_cell(s, x, y, cols[1], rh, row[1], BG)
        x += cols[1]
        table_cell(s, x, y, cols[2], rh, row[2], BG)
        x += cols[2]
        table_cell(s, x, y, cols[3], rh, row[3], BG)
        y += rh

    add_text(
        s,
        0.72,
        5.62,
        12.0,
        0.36,
        "适用场景：LRU（热点明显） / LFU（高频稳定） / FIFO（简单可预测）",
        size=12.2,
        color=TEXT_SUB,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        0.72,
        6.02,
        12.0,
        0.42,
        "一句话总结：统一抽象 + 策略可插拔 = 在性能、复杂度、可维护性之间取得最佳平衡。",
        size=15,
        bold=True,
        color=ACCENT_GREEN,
        align=PP_ALIGN.CENTER,
    )

    footer(s, 7)


# =========================
# Entry
# =========================
def build(output: Path):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)
    slide6(prs)
    slide7(prs)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def main():
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else (Path("docs") / "Spring-Cache-Demo.pptx")
    build(out)
    print(f"PPTX generated: {out.resolve()}")


if __name__ == "__main__":
    main()
